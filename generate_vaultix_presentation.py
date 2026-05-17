from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_NAME = "Vaultix_Soutenance.pptx"

COLORS = {
    "bg": RGBColor(5, 26, 51),
    "primary": RGBColor(10, 49, 97),
    "secondary": RGBColor(36, 101, 165),
    "accent": RGBColor(52, 225, 255),
    "text_light": RGBColor(214, 227, 240),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(11.9), Inches(0.9))
    p = title_box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = COLORS["text_light"]

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.3), Inches(11.6), Inches(0.7))
        sp = sub_box.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(17)
        sp.font.color.rgb = COLORS["accent"]


def add_bullets(slide, bullets, x=0.9, y=2.0, w=11.2, h=4.6, font_size=23):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS["text_light"]


def add_logo(slide, logo_path):
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), Inches(9.8), Inches(0.1), width=Inches(2.2), height=Inches(1.0))


def add_footer(slide, text="Vaultix — Secure Your Digital Life"):
    footer = slide.shapes.add_textbox(Inches(0), Inches(6.85), Inches(13.33), Inches(0.4))
    p = footer.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = COLORS["accent"]


def add_nav_buttons(slide, agenda_slide=None, prev_slide=None, next_slide=None):
    y = Inches(6.35)
    btn_w = Inches(1.1)
    btn_h = Inches(0.35)

    def _btn(x, label, target):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, btn_w, btn_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["secondary"]
        shape.line.color.rgb = COLORS["accent"]
        shape.text_frame.text = label
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        shape.text_frame.paragraphs[0].font.size = Pt(11)
        shape.text_frame.paragraphs[0].font.bold = True
        shape.text_frame.paragraphs[0].font.color.rgb = COLORS["text_light"]
        if target is not None:
            try:
                shape.click_action.target_slide = target
            except Exception:
                pass

    _btn(Inches(0.5), "Agenda", agenda_slide)
    _btn(Inches(1.75), "Prev", prev_slide)
    _btn(Inches(3.0), "Next", next_slide)


def main():
    base_dir = Path(__file__).resolve().parent
    logo_path = base_dir / "app" / "src" / "main" / "res" / "drawable" / "vaultix_logo.png"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1: Title
    s1 = prs.slides.add_slide(blank)
    add_bg(s1, COLORS["bg"])
    add_title(s1, "Vaultix", "Secure Your Digital Life")
    add_logo(s1, logo_path)
    add_bullets(
        s1,
        [
            "Android Password Manager — Kotlin + Firebase + AES",
            "Soutenance du mini-projet",
            "Authentification, sécurité, UX, architecture",
        ],
        y=2.4,
        font_size=24,
    )
    add_footer(s1)

    # Slide 2: Agenda
    s2 = prs.slides.add_slide(blank)
    add_bg(s2, COLORS["primary"])
    add_title(s2, "Agenda", "Plan de la présentation")
    agenda_items = [
        "1. Problématique et objectifs",
        "2. Architecture et technologies",
        "3. Sécurité (Master Password + chiffrement)",
        "4. Flux utilisateur + navigation",
        "5. Démo fonctionnalités",
        "6. Difficultés, solutions et perspectives",
    ]
    add_bullets(s2, agenda_items, y=1.95, font_size=22)
    add_logo(s2, logo_path)
    add_footer(s2)

    # Content slides
    slides = []

    content = [
        (
            "Problématique",
            "Pourquoi Vaultix ?",
            [
                "Multiplication des comptes et mots de passe",
                "Risque élevé de réutilisation de secrets faibles",
                "Besoin: centraliser + chiffrer + simplifier l'usage",
                "Vaultix répond avec un coffre chiffré côté client",
            ],
        ),
        (
            "Objectifs",
            "Ce que le projet devait livrer",
            [
                "Authentification Firebase (email + Google)",
                "Master Password demandé à chaque session",
                "Chiffrement AES-GCM de toutes les entrées",
                "UI claire: Welcome, Login, Master, Home, Add",
            ],
        ),
        (
            "Architecture",
            "Vue technique du projet",
            [
                "Single-Activity: MainActivity orchestre tout",
                "Fragments: Welcome/Login/Master/Home/Add",
                "Repository: FirebaseRepository (Auth + Firestore)",
                "Crypto: PasswordCrypto (PBKDF2 + AES-GCM)",
            ],
        ),
        (
            "Sécurité",
            "Master Password et dérivation de clé",
            [
                "Salt unique par utilisateur stocké en base",
                "KDF: PBKDF2WithHmacSHA256 (120000 itérations)",
                "Clé AES 256 bits dérivée à chaque session",
                "Aucune donnée sensible en clair dans Firestore",
            ],
        ),
        (
            "Validation Master",
            "1 Master Password effectif par utilisateur",
            [
                "Création d'un verifier chiffré au premier unlock",
                "À chaque session: déchiffrement du verifier",
                "Mismatch -> accès refusé",
                "Permet de bloquer les faux master passwords",
            ],
        ),
        (
            "Flux Utilisateur",
            "De l'ouverture à la gestion des mots de passe",
            [
                "Welcome -> Login/ Master selon session Firebase",
                "Master validé -> Home",
                "Home -> Add via FAB",
                "Drawer: Profile/Settings/Logout/Other",
            ],
        ),
        (
            "Fonctionnalités UX",
            "Ce qui a été ajouté en plus",
            [
                "Génération password fort (A-Za-z0-9)",
                "Show/Hide password dans la liste",
                "Copy to clipboard pour chaque entrée",
                "Welcome/Loading avec branding",
            ],
        ),
        (
            "Google Sign-In",
            "Intégration complète",
            [
                "Account picker Google + récupération ID token",
                "Firebase signInWithCredential(token)",
                "Nécessite SHA-1/SHA-256 + google-services.json à jour",
                "Gestion d'erreurs claire côté app",
            ],
        ),
        (
            "Démonstration",
            "Script live",
            [
                "1) Ouvrir app -> Welcome",
                "2) Login (email ou Google)",
                "3) Entrer Master Password",
                "4) Ajouter entrée + générer password",
                "5) Show/Hide + Copy + Logout",
            ],
        ),
        (
            "Difficultés et Solutions",
            "Ce qui a été résolu",
            [
                "CONFIGURATION_NOT_FOUND -> Auth Firebase setup",
                "PERMISSION_DENIED -> règles Firestore",
                "Google no-op -> SHA/OAuth/Web client id",
                "Navigation back/drawer -> logique MainActivity corrigée",
            ],
        ),
        (
            "Conclusion",
            "Bilan du mini-projet",
            [
                "Application fonctionnelle, sécurisée et démontrable",
                "Architecture propre et évolutive",
                "Base solide pour biométrie, edit/delete, offline cache",
                "Excellent premier projet Android complet",
            ],
        ),
    ]

    for title, subtitle, bullets in content:
        slide = prs.slides.add_slide(blank)
        add_bg(slide, COLORS["bg"])
        add_title(slide, title, subtitle)
        add_bullets(slide, bullets)
        add_logo(slide, logo_path)
        add_footer(slide)
        slides.append(slide)

    # Interactive agenda buttons -> jump to each section slide
    cards_start_y = 4.95
    card_w = Inches(3.9)
    card_h = Inches(0.62)
    x_positions = [Inches(0.9), Inches(4.75), Inches(8.6)]

    targets = [slides[0], slides[2], slides[4], slides[6], slides[8], slides[10]]
    labels = [
        "Problématique",
        "Architecture",
        "Master Security",
        "Fonctionnalités",
        "Démo",
        "Conclusion",
    ]

    for idx, (label, target) in enumerate(zip(labels, targets)):
        row = idx // 3
        col = idx % 3
        shape = s2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_positions[col],
            cards_start_y + Inches(row * 0.75),
            card_w,
            card_h,
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["secondary"]
        shape.line.color.rgb = COLORS["accent"]
        shape.text_frame.text = label
        p = shape.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLORS["white"]
        try:
            shape.click_action.target_slide = target
        except Exception:
            pass

    # Add nav buttons to all slides except title
    all_slides = list(prs.slides)
    for i, slide in enumerate(all_slides):
        if i == 0:
            continue
        prev_slide = all_slides[i - 1] if i - 1 >= 0 else None
        next_slide = all_slides[i + 1] if i + 1 < len(all_slides) else None
        add_nav_buttons(slide, agenda_slide=s2, prev_slide=prev_slide, next_slide=next_slide)

    output_path = base_dir / OUTPUT_NAME
    prs.save(output_path)
    print(f"Presentation generated: {output_path}")


if __name__ == "__main__":
    main()
